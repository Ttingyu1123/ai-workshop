from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "ppt" / "ai-workshop-demo.pptx"

W, H = Inches(13.333), Inches(7.5)

PRIMARY = RGBColor(0x49, 0xB6, 0xE5)
SECONDARY = RGBColor(0x26, 0x3D, 0x5B)
TEXT = RGBColor(0x11, 0x18, 0x27)
TEXT_LIGHT = RGBColor(0x4B, 0x55, 0x63)
PAPER = RGBColor(0xFE, 0xFC, 0xF3)
PAPER_LINE = RGBColor(0xE8, 0xE4, 0xD9)
YELLOW = RGBColor(0xFD, 0xE6, 0x8A)
GREEN = RGBColor(0xBB, 0xF7, 0xD0)
BLUE = RGBColor(0xBA, 0xE6, 0xFD)
PINK = RGBColor(0xFB, 0xCF, 0xE8)
PURPLE = RGBColor(0xDD, 0xD6, 0xFE)
DANGER_BG = RGBColor(0xFF, 0xEB, 0xEE)
DANGER = RGBColor(0xDC, 0x26, 0x26)
SUCCESS_BG = RGBColor(0xE8, 0xF5, 0xE9)
SUCCESS = RGBColor(0x16, 0xA3, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

TITLE_FONT = "LXGW WenKai TC"
HAND_FONT = "Caveat"
BODY_FONT = "LXGW WenKai TC"
MONO_FONT = "JetBrains Mono"


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def set_line(shape, color=TEXT, width=1.8):
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def add_textbox(slide, text, x, y, w, h, font_size=22, color=TEXT, bold=False,
                font=BODY_FONT, align=None, valign=MSO_ANCHOR.TOP, line_spacing=1.1):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.text = text
    p.line_spacing = line_spacing
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_paragraphs(slide, lines, x, y, w, h, font_size=19, color=TEXT, bullet=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_after = Pt(7)
        p.line_spacing = 1.15
        if bullet:
            p.text = line
            p.font.name = BODY_FONT
        run = p.runs[0]
        run.font.name = BODY_FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
    return box


def add_bg(slide):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = PAPER
    y = 0.36
    while y < 7.5:
        line = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(0), Inches(y), Inches(13.333), Inches(y)
        )
        line.line.color.rgb = PAPER_LINE
        line.line.width = Pt(0.5)
        y += 0.36


def add_footer(slide, page):
    add_textbox(
        slide,
        "衛生福利部台中醫院 感染科 曾婷玉醫師 · AI Workshop",
        0.55,
        7.12,
        8.2,
        0.25,
        font_size=8.5,
        color=TEXT_LIGHT,
        font=BODY_FONT,
    )
    add_textbox(
        slide,
        f"{page:02d}",
        12.35,
        7.05,
        0.45,
        0.28,
        font_size=9.5,
        color=TEXT_LIGHT,
        font=MONO_FONT,
        align=PP_ALIGN.RIGHT,
    )


def add_badge(slide, text, x, y, w=1.5, color=YELLOW, rotate=-3):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.38)
    )
    set_fill(shape, color)
    set_line(shape, TEXT, 1.2)
    shape.rotation = rotate
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    run = p.runs[0]
    run.font.name = HAND_FONT
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = TEXT
    return shape


def add_card(slide, x, y, w, h, fill=WHITE, rotate=0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    set_fill(shape, fill)
    set_line(shape, TEXT, 2)
    shape.rotation = rotate
    return shape


def add_card_text(
    slide,
    title,
    body,
    x,
    y,
    w,
    h,
    fill=WHITE,
    tag=None,
    tag_color=YELLOW,
    rotate=0,
    body_font_size=14.5,
):
    add_card(slide, x, y, w, h, fill, rotate)
    if tag:
        add_badge(slide, tag, x + 0.18, y + 0.2, w=1.05, color=tag_color, rotate=-2)
        title_y = y + 0.72
    else:
        title_y = y + 0.32
    add_textbox(slide, title, x + 0.28, title_y, w - 0.55, 0.35, font_size=18, bold=True, font=BODY_FONT)
    body_y = title_y + 0.42
    body_h = max(0.35, y + h - body_y - 0.18)
    add_textbox(
        slide,
        body,
        x + 0.28,
        body_y,
        w - 0.55,
        body_h,
        font_size=body_font_size,
        color=TEXT_LIGHT,
        font=BODY_FONT,
        line_spacing=0.95,
    )


def add_title(slide, title, subtitle=None, badge=None, section_color=YELLOW):
    if badge:
        add_badge(slide, badge, 0.68, 0.48, w=1.45, color=section_color)
    add_textbox(slide, title, 0.68, 0.94, 8.3, 0.75, font_size=32, bold=True, font=TITLE_FONT)
    if subtitle:
        add_textbox(slide, subtitle, 0.7, 1.67, 8.6, 0.35, font_size=15.5, color=TEXT_LIGHT, font=BODY_FONT)
    wave_y = 2.1
    for i in range(10):
        x1 = 0.72 + i * 0.16
        y1 = wave_y + (0.04 if i % 2 == 0 else -0.02)
        x2 = x1 + 0.14
        y2 = wave_y + (-0.02 if i % 2 == 0 else 0.04)
        seg = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
        )
        seg.line.color.rgb = TEXT
        seg.line.width = Pt(1.2)


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_badge(slide, "Pilot Deck", 0.75, 0.65, w=1.6, color=YELLOW)
    add_textbox(slide, "AI 入門與創作工作坊", 0.75, 1.38, 6.35, 0.9, font_size=35, bold=True, font=TITLE_FONT)
    add_textbox(slide, "給醫療人員的第一堂 AI 課", 0.78, 2.24, 5.7, 0.4, font_size=18, color=TEXT_LIGHT)
    add_card(slide, 0.78, 3.08, 5.55, 1.75, fill=WHITE, rotate=-0.5)
    add_textbox(slide, "今天先建立三件事", 1.08, 3.34, 4.7, 0.3, font_size=17, bold=True)
    add_paragraphs(
        slide,
        ["知道 AI 能做什麼，也知道不能信什麼", "會寫出有邊界的 Prompt", "先守住病人隱私與專業責任"],
        1.12,
        3.78,
        4.72,
        0.8,
        font_size=14.2,
    )
    add_card(slide, 7.65, 0.82, 4.65, 4.65, fill=WHITE, rotate=1.2)
    slide.shapes.add_picture(str(ROOT / "og-image.png"), Inches(7.87), Inches(1.04), width=Inches(4.2), height=Inches(4.2))
    add_textbox(slide, "Doodle style from GitHub Pages", 7.95, 5.68, 4.3, 0.35, font_size=13.5, color=TEXT_LIGHT, font=HAND_FONT, align=PP_ALIGN.CENTER)
    add_footer(slide, 1)


def slide_map(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "課程地圖", "從安全使用開始，再進到文字、圖像、影片與 Agent。", badge="Workshop")
    items = [
        ("基礎必修", "AI 基礎觀念\n文字應用\n安全與隱私", GREEN, "Part 1-3"),
        ("創作應用", "圖像生成\n影片製作\nGoogle Flow", BLUE, "Part 4-6"),
        ("進階實戰", "AI Agent\n自動化任務\n工具協作", PINK, "Part 7"),
        ("番外工具", "NotebookLM\n語音轉文字\nCanva / CapCut / GitHub Pages", PURPLE, "Extras"),
    ]
    for i, (title, body, color, tag) in enumerate(items):
        x = 0.8 + (i % 2) * 6.0
        y = 2.4 + (i // 2) * 2.22
        add_card_text(
            slide,
            title,
            body,
            x,
            y,
            5.35,
            1.9,
            fill=WHITE,
            tag=tag,
            tag_color=color,
            rotate=(-0.5 if i % 2 == 0 else 0.4),
            body_font_size=13.2,
        )
    add_footer(slide, 2)


def slide_ai_concept(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "AI 到底在幹嘛", "用醫療人員熟悉的方式理解 LLM。", badge="Part 1")
    add_card(slide, 0.85, 2.55, 5.55, 2.55, fill=GREEN, rotate=-0.3)
    add_textbox(slide, "醫學類比", 1.18, 2.83, 4.6, 0.32, font_size=17, bold=True, color=SUCCESS)
    add_textbox(
        slide,
        "像一個讀完所有教科書的住院醫師：知識量驚人，但沒有臨床經驗，偶爾會很有自信地說錯話。",
        1.18,
        3.28,
        4.75,
        1.1,
        font_size=18,
        font=BODY_FONT,
    )
    add_textbox(slide, "LLM = 預測下一個字", 7.1, 2.6, 5.1, 0.55, font_size=28, bold=True, color=SECONDARY)
    add_paragraphs(
        slide,
        ["不是資料庫查詢，而是根據上下文生成", "回答流暢不等於已經查證", "越明確的任務，越容易得到可用結果"],
        7.25,
        3.38,
        4.85,
        1.2,
        font_size=17.5,
    )
    add_badge(slide, "keyword", 7.22, 4.92, w=1.25, color=YELLOW)
    add_textbox(slide, "Prompt 是你開給 AI 的「醫囑」。", 8.62, 4.95, 3.1, 0.35, font_size=17, bold=True)
    add_footer(slide, 3)


def slide_safety(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "醫療場景先守住紅線", "能不能用 AI，第一題永遠是資料可不可以貼。", badge="Part 3", section_color=PINK)
    add_card(slide, 0.82, 2.45, 5.5, 3.35, fill=DANGER_BG, rotate=-0.4)
    add_textbox(slide, "不要這樣貼", 1.15, 2.76, 4.5, 0.35, font_size=20, bold=True, color=DANGER)
    add_textbox(
        slide,
        "王小明，78 歲男性，病歷號 A123456789，住台北市信義區，因肺癌合併肋膜積水就診...",
        1.15,
        3.28,
        4.65,
        1.15,
        font_size=17,
    )
    add_textbox(slide, "直接識別 + 間接識別組合，風險會疊加。", 1.15, 4.75, 4.65, 0.38, font_size=15, color=TEXT_LIGHT)
    add_card(slide, 7.05, 2.45, 5.5, 3.35, fill=SUCCESS_BG, rotate=0.4)
    add_textbox(slide, "改成這樣", 7.38, 2.76, 4.5, 0.35, font_size=20, bold=True, color=SUCCESS)
    add_textbox(
        slide,
        "一位高齡男性病患，因肺癌合併肋膜積水於胸腔科就診，同時有糖尿病與慢性病毒性肝炎病史。",
        7.38,
        3.28,
        4.65,
        1.15,
        font_size=17,
    )
    add_textbox(slide, "保留臨床意義，移除可回推身分的細節。", 7.38, 4.75, 4.65, 0.38, font_size=15, color=TEXT_LIGHT)
    add_footer(slide, 4)


def slide_prompt(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "一個好 Prompt 的骨架", "不是問得很聰明，而是交代得很清楚。", badge="Part 2", section_color=BLUE)
    cards = [
        ("角色", "你是一位感染科醫師，正在為護理師做教育訓練", GREEN),
        ("任務", "請整理成 5 分鐘可講完的衛教重點", BLUE),
        ("格式", "用表格呈現，欄位含重點、說法、注意事項", YELLOW),
        ("限制", "不要使用病人個資；避免診斷建議與藥物劑量", PINK),
    ]
    for i, (title, body, color) in enumerate(cards):
        x = 0.78 + i * 3.15
        add_card_text(slide, title, body, x, 2.35, 2.65, 1.72, fill=WHITE, tag=str(i + 1), tag_color=color, rotate=(-0.5 if i % 2 else 0.4))
    add_card(slide, 1.18, 4.7, 10.95, 1.18, fill=WHITE)
    top = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(1.18), Inches(4.7), Inches(10.95), Inches(0.1))
    set_fill(top, PRIMARY)
    top.line.fill.background()
    add_textbox(
        slide,
        "請幫我把以下出院衛教改寫成國中生能懂的版本，分成「為什麼重要、怎麼做、什麼時候回診」三段；不要新增未提供的藥物劑量。",
        1.5,
        4.98,
        10.25,
        0.54,
        font_size=17,
        font=BODY_FONT,
    )
    add_footer(slide, 5)


def slide_activity(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_title(slide, "現在就做：10 分鐘實作", "把一段真實工作改成可安全貼給 AI 的任務。", badge="Exercise", section_color=YELLOW)
    add_card(slide, 0.8, 2.4, 7.0, 3.45, fill=WHITE, rotate=-0.3)
    add_textbox(slide, "任務流程", 1.1, 2.72, 5.8, 0.35, font_size=20, bold=True)
    add_paragraphs(
        slide,
        ["1. 選一段不含個資的衛教或行政文字", "2. 寫出角色、任務、格式、限制", "3. 產出後標出 2 個需要人工核對的地方", "4. 分享：哪一句 Prompt 改完最有差？"],
        1.15,
        3.22,
        6.2,
        1.65,
        font_size=17,
    )
    add_card(slide, 8.55, 2.32, 3.35, 3.35, fill=WHITE, rotate=0.7)
    slide.shapes.add_picture(str(ROOT / "qr-site.png"), Inches(8.88), Inches(2.66), width=Inches(2.7), height=Inches(2.7))
    add_textbox(slide, "講義與工具頁", 8.9, 5.82, 2.6, 0.3, font_size=14.5, font=HAND_FONT, color=TEXT_LIGHT, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)


def build():
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    slide_cover(prs)
    slide_map(prs)
    slide_ai_concept(prs)
    slide_safety(prs)
    slide_prompt(prs)
    slide_activity(prs)
    prs.save(OUT)


if __name__ == "__main__":
    build()
    print(OUT)
